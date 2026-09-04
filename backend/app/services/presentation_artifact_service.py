"""Deterministic in-memory PowerPoint generation from trusted presentation data."""

import json
import posixpath
import re
from datetime import datetime
from decimal import Decimal
from hashlib import sha256
from io import BytesIO
from math import ceil
from pathlib import PurePosixPath
from typing import cast
from urllib.parse import urlsplit
from xml.etree import ElementTree
from zipfile import BadZipFile, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationDocument
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt
from pydantic import ValidationError

from backend.app.schemas.presentation import AnalyticalPresentationResult
from backend.app.schemas.presentation_artifact import (
    MAX_PRESENTATION_ARTIFACT_BYTES,
    MAX_PRESENTATION_SLIDES,
    PPTX_MIME_TYPE,
    PresentationArtifact,
)
from backend.app.schemas.visualization import (
    BarVisualizationSpec,
    KPIVisualizationSpec,
    LineVisualizationSpec,
    TableVisualizationSpec,
)

TABLE_ROWS_PER_SLIDE = 12
SERVER_GENERATED_TITLE = "Analytical Presentation"
FIXED_DOCUMENT_TIMESTAMP = datetime(2000, 1, 1)


def _rgb(red: int, green: int, blue: int) -> RGBColor:
    return RGBColor(red, green, blue)  # type: ignore[no-untyped-call]


BACKGROUND_COLOR = _rgb(248, 250, 252)
PRIMARY_COLOR = _rgb(15, 76, 129)
ACCENT_COLOR = _rgb(14, 116, 144)
NEGATIVE_COLOR = _rgb(185, 28, 28)
TEXT_COLOR = _rgb(15, 23, 42)
MUTED_COLOR = _rgb(71, 85, 105)
WHITE_COLOR = _rgb(255, 255, 255)
GRID_COLOR = _rgb(203, 213, 225)


class PresentationArtifactError(RuntimeError):
    """Base error raised by deterministic presentation artifact generation."""


class PresentationArtifactInputError(PresentationArtifactError):
    """Raised when the trusted presentation source is inconsistent or unsupported."""


class PresentationArtifactSizeError(PresentationArtifactError):
    """Raised when the generated artifact exceeds its controlled byte limit."""


def _format_decimal(value: Decimal, unit: str) -> str:
    rendered = format(value, "f")

    if unit == "brl":
        return f"R$ {rendered}"

    if unit == "percentage":
        return f"{rendered}%"

    return rendered


def _add_text_box(
    slide: Slide,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    color: RGBColor = TEXT_COLOR,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> Shape:
    shape = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = alignment
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return shape


def _set_background(slide: Slide) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = BACKGROUND_COLOR


def _add_slide_header(slide: Slide, title: str, subtitle: str) -> None:
    _set_background(slide)
    _add_text_box(
        slide,
        title,
        left=0.65,
        top=0.35,
        width=12.0,
        height=0.55,
        font_size=26,
        color=PRIMARY_COLOR,
        bold=True,
    )
    _add_text_box(
        slide,
        subtitle,
        left=0.68,
        top=0.92,
        width=11.9,
        height=0.32,
        font_size=11,
        color=MUTED_COLOR,
    )


def _new_slide(presentation: PresentationDocument) -> Slide:
    slide = cast(Slide, presentation.slides.add_slide(presentation.slide_layouts[6]))
    _set_background(slide)
    return slide


def _semantic_artifact_id(result: AnalyticalPresentationResult) -> str:
    safe_payload = {
        "presentation_version": result.presentation_version,
        "source_row_count": result.source_row_count,
        "visualizations": result.visualizations.model_dump(mode="json"),
        "insights": result.insights.model_dump(
            mode="json",
            include={
                "insight_version",
                "analytics_version",
                "visualization_version",
                "execution_version",
                "semantic_version",
                "catalog_version",
                "source_row_count",
                "summary",
                "claims",
            },
        ),
    }
    canonical_payload = json.dumps(
        safe_payload,
        ensure_ascii=True,
        separators=(",", ":"),
        sort_keys=True,
    ).encode("utf-8")
    return sha256(canonical_payload).hexdigest()[:24]


def _table_slide_count(specification: TableVisualizationSpec) -> int:
    return ceil(len(specification.rows) / TABLE_ROWS_PER_SLIDE)


def _required_slide_count(result: AnalyticalPresentationResult) -> int:
    visualization_slides = sum(
        _table_slide_count(specification)
        if isinstance(specification, TableVisualizationSpec)
        else 1
        for specification in result.visualizations.specifications
    )
    return 3 + visualization_slides + len(result.insights.claims)


def _validate_source(result: AnalyticalPresentationResult) -> int:
    if not isinstance(result, AnalyticalPresentationResult):
        raise PresentationArtifactInputError("Presentation artifact source is invalid")

    visualizations = result.visualizations
    insights = result.insights

    if result.presentation_status != "generated":
        raise PresentationArtifactInputError("Presentation artifact source is not generated")

    if result.source_row_count < 1:
        raise PresentationArtifactInputError("Presentation artifact requires source rows")

    if visualizations.visualization_status != "specified" or not visualizations.deterministic:
        raise PresentationArtifactInputError("Presentation visualizations are not trusted")

    if insights.insight_status != "generated" or not insights.grounded:
        raise PresentationArtifactInputError("Presentation insights are not grounded")

    if insights.calculated_by_llm:
        raise PresentationArtifactInputError(
            "Presentation insights cannot contain LLM calculations"
        )

    if not visualizations.specifications:
        raise PresentationArtifactInputError("Presentation artifact requires visualizations")

    if not insights.claims:
        raise PresentationArtifactInputError("Presentation artifact requires grounded claims")

    matching_fields = (
        "visualization_version",
        "analytics_version",
        "execution_version",
        "semantic_version",
        "catalog_version",
        "source_row_count",
    )

    for field_name in matching_fields:
        if getattr(visualizations, field_name) != getattr(insights, field_name):
            raise PresentationArtifactInputError(
                f"Presentation artifact source mismatch: {field_name}"
            )

    if result.source_row_count != visualizations.source_row_count:
        raise PresentationArtifactInputError("Presentation artifact row count is inconsistent")

    slide_count = _required_slide_count(result)

    if slide_count > MAX_PRESENTATION_SLIDES:
        raise PresentationArtifactInputError("Presentation artifact exceeds the slide limit")

    return slide_count


def _add_title_slide(
    presentation: PresentationDocument,
    artifact_id: str,
    source_row_count: int,
) -> None:
    slide = _new_slide(presentation)
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(1.2),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = PRIMARY_COLOR
    banner.line.fill.background()
    _add_text_box(
        slide,
        SERVER_GENERATED_TITLE,
        left=0.85,
        top=2.0,
        width=11.65,
        height=1.0,
        font_size=34,
        color=PRIMARY_COLOR,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide,
        "Deterministic analytics, visualizations, and grounded insights",
        left=1.35,
        top=3.05,
        width=10.65,
        height=0.7,
        font_size=18,
        color=MUTED_COLOR,
        alignment=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide,
        f"Artifact {artifact_id} · Source rows {source_row_count}",
        left=1.5,
        top=5.65,
        width=10.3,
        height=0.4,
        font_size=12,
        color=MUTED_COLOR,
        alignment=PP_ALIGN.CENTER,
    )


def _add_kpi_slide(
    presentation: PresentationDocument,
    specification: KPIVisualizationSpec,
) -> None:
    slide = _new_slide(presentation)
    _add_slide_header(slide, specification.title, specification.spec_id)
    _add_text_box(
        slide,
        _format_decimal(specification.value, specification.unit),
        left=1.0,
        top=1.65,
        width=11.3,
        height=1.35,
        font_size=40,
        color=ACCENT_COLOR,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    details = (
        ("Aggregation", str(specification.aggregation)),
        ("Count", str(specification.value_count)),
        ("Total", _format_decimal(specification.total, specification.unit)),
        ("Average", _format_decimal(specification.average, specification.unit)),
        ("Minimum", _format_decimal(specification.minimum, specification.unit)),
        ("Maximum", _format_decimal(specification.maximum, specification.unit)),
    )

    for index, (label, value) in enumerate(details):
        column = index % 3
        row = index // 3
        left = 0.8 + (column * 4.15)
        top = 3.5 + (row * 1.2)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(3.75),
            Inches(0.9),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE_COLOR
        card.line.color.rgb = GRID_COLOR
        _add_text_box(
            slide,
            label,
            left=left + 0.2,
            top=top + 0.08,
            width=1.2,
            height=0.28,
            font_size=10,
            color=MUTED_COLOR,
        )
        _add_text_box(
            slide,
            value,
            left=left + 0.2,
            top=top + 0.37,
            width=3.35,
            height=0.35,
            font_size=15,
            color=TEXT_COLOR,
            bold=True,
        )


def _add_table_slides(
    presentation: PresentationDocument,
    specification: TableVisualizationSpec,
) -> None:
    page_count = _table_slide_count(specification)

    for page_index in range(page_count):
        start = page_index * TABLE_ROWS_PER_SLIDE
        rows = specification.rows[start : start + TABLE_ROWS_PER_SLIDE]
        slide = _new_slide(presentation)
        page_label = f"Page {page_index + 1} of {page_count}"
        _add_slide_header(
            slide,
            specification.title,
            f"{specification.spec_id} · {page_label}",
        )
        graphic_frame = slide.shapes.add_table(
            len(rows) + 1,
            4,
            Inches(0.65),
            Inches(1.45),
            Inches(12.0),
            Inches(5.2),
        )
        table = graphic_frame.table
        headings = ("Position", specification.dimension_name, "Value", "Share")

        for column, heading in enumerate(headings):
            cell = table.cell(0, column)
            cell.text = heading
            cell.fill.solid()  # type: ignore[no-untyped-call]
            cell.fill.fore_color.rgb = PRIMARY_COLOR
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE_COLOR

        for row_index, row in enumerate(rows, start=1):
            values = (
                str(row.position),
                row.label,
                _format_decimal(row.value, specification.unit),
                "" if row.share_percent is None else f"{format(row.share_percent, 'f')}%",
            )

            for column, value in enumerate(values):
                cell = table.cell(row_index, column)
                cell.text = value
                cell.fill.solid()  # type: ignore[no-untyped-call]
                cell.fill.fore_color.rgb = WHITE_COLOR
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.name = "Aptos"
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = TEXT_COLOR

        _add_text_box(
            slide,
            "Full ranking total: "
            f"{_format_decimal(specification.ranking_total, specification.unit)}",
            left=0.72,
            top=6.72,
            width=11.8,
            height=0.3,
            font_size=10,
            color=MUTED_COLOR,
        )


def _add_bar_slide(
    presentation: PresentationDocument,
    specification: BarVisualizationSpec,
) -> None:
    slide = _new_slide(presentation)
    _add_slide_header(slide, specification.title, specification.spec_id)
    maximum = max(abs(item.value) for item in specification.items)
    maximum = maximum if maximum != 0 else Decimal("1")
    available_width = 7.0
    row_height = min(0.25, 4.8 / len(specification.items))
    row_step = 4.9 / len(specification.items)

    for index, item in enumerate(specification.items):
        top = 1.45 + (index * row_step)
        width = float(abs(item.value) / maximum) * available_width
        width = max(width, 0.04)
        _add_text_box(
            slide,
            item.label,
            left=0.7,
            top=top,
            width=2.25,
            height=row_height,
            font_size=9,
            color=TEXT_COLOR,
        )
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3.0),
            Inches(top),
            Inches(width),
            Inches(row_height),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT_COLOR if item.value >= 0 else NEGATIVE_COLOR
        bar.line.fill.background()
        _add_text_box(
            slide,
            _format_decimal(item.value, specification.unit),
            left=10.2,
            top=top,
            width=2.35,
            height=row_height,
            font_size=9,
            color=TEXT_COLOR,
            alignment=PP_ALIGN.RIGHT,
        )

    _add_text_box(
        slide,
        f"Full ranking total: {_format_decimal(specification.ranking_total, specification.unit)}",
        left=0.72,
        top=6.72,
        width=11.8,
        height=0.3,
        font_size=10,
        color=MUTED_COLOR,
    )


def _line_coordinates(
    specification: LineVisualizationSpec,
) -> tuple[tuple[Emu, Emu], ...]:
    left = Inches(1.0)
    top = Inches(1.65)
    width = Inches(11.2)
    height = Inches(4.6)
    values = tuple(point.value for point in specification.points)
    minimum = min(values)
    maximum = max(values)
    span = maximum - minimum
    count = len(values)
    coordinates: list[tuple[Emu, Emu]] = []

    for index, value in enumerate(values):
        x_ratio = 0.5 if count == 1 else index / (count - 1)
        y_ratio = 0.5 if span == 0 else float((maximum - value) / span)
        x = Emu(int(left + (width * x_ratio)))
        y = Emu(int(top + (height * y_ratio)))
        coordinates.append((x, y))

    return tuple(coordinates)


def _add_line_slide(
    presentation: PresentationDocument,
    specification: LineVisualizationSpec,
) -> None:
    slide = _new_slide(presentation)
    _add_slide_header(slide, specification.title, specification.spec_id)
    coordinates = _line_coordinates(specification)

    for first, second in zip(coordinates, coordinates[1:], strict=False):
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            first[0],
            first[1],
            second[0],
            second[1],
        )
        connector.line.color.rgb = ACCENT_COLOR
        connector.line.width = Pt(2)

    label_interval = max(1, len(coordinates) // 8)

    for index, ((x, y), point) in enumerate(zip(coordinates, specification.points, strict=True)):
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(x - Inches(0.045)),
            Emu(y - Inches(0.045)),
            Inches(0.09),
            Inches(0.09),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = PRIMARY_COLOR
        marker.line.fill.background()

        if index % label_interval == 0 or index == len(coordinates) - 1:
            _add_text_box(
                slide,
                f"{point.label}\n{_format_decimal(point.value, specification.unit)}",
                left=float(x) / 914_400 - 0.45,
                top=6.35,
                width=0.9,
                height=0.55,
                font_size=8,
                color=MUTED_COLOR,
                alignment=PP_ALIGN.CENTER,
            )


def _evidence_label(evidence_type: str, metric_name: str | None, spec_id: str | None) -> str:
    if evidence_type == "visualization":
        if spec_id is None:
            raise PresentationArtifactInputError("Visualization evidence is incomplete")
        return f"Visualization evidence: {spec_id}"

    if metric_name is None:
        raise PresentationArtifactInputError("Metric evidence is incomplete")
    return f"Metric evidence: {metric_name}"


def _add_insight_summary_slide(
    presentation: PresentationDocument,
    result: AnalyticalPresentationResult,
) -> None:
    slide = _new_slide(presentation)
    _add_slide_header(slide, "Grounded insight summary", result.insights.insight_version)
    _add_text_box(
        slide,
        result.insights.summary,
        left=1.0,
        top=1.7,
        width=11.3,
        height=3.2,
        font_size=22,
        color=TEXT_COLOR,
    )
    _add_text_box(
        slide,
        f"Grounded claims: {len(result.insights.claims)}",
        left=1.0,
        top=5.35,
        width=11.3,
        height=0.45,
        font_size=13,
        color=MUTED_COLOR,
    )


def _add_claim_slides(
    presentation: PresentationDocument,
    result: AnalyticalPresentationResult,
) -> None:
    for index, claim in enumerate(result.insights.claims, start=1):
        slide = _new_slide(presentation)
        _add_slide_header(
            slide,
            f"Grounded claim {index}",
            claim.claim_id,
        )
        _add_text_box(
            slide,
            claim.text,
            left=0.9,
            top=1.45,
            width=11.55,
            height=2.5,
            font_size=20,
            color=TEXT_COLOR,
        )
        evidence_lines = tuple(
            _evidence_label(
                reference.evidence_type,
                reference.metric_name,
                reference.specification_id,
            )
            for reference in claim.evidence
        )
        _add_text_box(
            slide,
            "\n".join(evidence_lines),
            left=1.0,
            top=4.25,
            width=11.3,
            height=1.65,
            font_size=14,
            color=PRIMARY_COLOR,
        )


def _add_provenance_slide(
    presentation: PresentationDocument,
    result: AnalyticalPresentationResult,
    artifact_id: str,
) -> None:
    slide = _new_slide(presentation)
    _add_slide_header(slide, "Source provenance", artifact_id)
    provenance = (
        ("Presentation version", result.presentation_version),
        ("Visualization version", result.visualizations.visualization_version),
        ("Insight version", result.insights.insight_version),
        ("Analytics version", result.visualizations.analytics_version),
        ("Execution version", result.visualizations.execution_version),
        ("Semantic version", result.visualizations.semantic_version),
        ("Catalog version", result.visualizations.catalog_version),
        ("Source row count", str(result.source_row_count)),
    )

    for index, (label, value) in enumerate(provenance):
        row = index // 2
        column = index % 2
        left = 0.9 + (column * 6.0)
        top = 1.5 + (row * 1.15)
        _add_text_box(
            slide,
            label,
            left=left,
            top=top,
            width=2.3,
            height=0.3,
            font_size=11,
            color=MUTED_COLOR,
        )
        _add_text_box(
            slide,
            value,
            left=left,
            top=top + 0.32,
            width=5.4,
            height=0.42,
            font_size=15,
            color=TEXT_COLOR,
            bold=True,
        )


def _build_presentation(
    result: AnalyticalPresentationResult,
    artifact_id: str,
) -> PresentationDocument:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    properties = presentation.core_properties
    properties.title = SERVER_GENERATED_TITLE
    properties.subject = "Deterministic analytical presentation"
    properties.author = "AI SQL Data Analyst"
    properties.last_modified_by = "AI SQL Data Analyst"
    properties.comments = "Generated from validated deterministic specifications"
    properties.created = FIXED_DOCUMENT_TIMESTAMP
    properties.modified = FIXED_DOCUMENT_TIMESTAMP
    _add_title_slide(presentation, artifact_id, result.source_row_count)

    for specification in result.visualizations.specifications:
        if isinstance(specification, KPIVisualizationSpec):
            _add_kpi_slide(presentation, specification)
        elif isinstance(specification, TableVisualizationSpec):
            _add_table_slides(presentation, specification)
        elif isinstance(specification, BarVisualizationSpec):
            _add_bar_slide(presentation, specification)
        elif isinstance(specification, LineVisualizationSpec):
            _add_line_slide(presentation, specification)
        else:
            raise PresentationArtifactInputError("Presentation visualization type is unsupported")

    _add_insight_summary_slide(presentation, result)
    _add_claim_slides(presentation, result)
    _add_provenance_slide(presentation, result, artifact_id)
    return presentation


_DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
_OFFICE_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_PPT_CT = "application/vnd.openxmlformats-officedocument.presentationml."
_OOXML_MEMBER = re.compile(
    r"(?:\[Content_Types\]\.xml|_rels/\.rels|"
    r"docProps/(?:core\.xml|app\.xml|thumbnail\.jpeg)|"
    r"ppt/(?:presentation|presProps|viewProps|tableStyles)\.xml|"
    r"ppt/_rels/presentation\.xml\.rels|"
    r"ppt/theme/theme[1-9][0-9]*\.xml|"
    r"ppt/printerSettings/printerSettings1\.bin|"
    r"ppt/slides/(?:slide[1-9][0-9]*\.xml|_rels/slide[1-9][0-9]*\.xml\.rels)|"
    r"ppt/slideMasters/(?:slideMaster[1-9][0-9]*\.xml|_rels/slideMaster[1-9][0-9]*\.xml\.rels)|"
    r"ppt/slideLayouts/(?:slideLayout[1-9][0-9]*\.xml|_rels/slideLayout[1-9][0-9]*\.xml\.rels))"
)
_PROHIBITED_XML_ELEMENTS = frozenset(
    {
        "oleobj",
        "oleobject",
        "control",
        "controls",
        "chart",
        "externaldata",
        "hlinkclick",
        "hlinkmouseover",
        "audio",
        "video",
        "audiofile",
        "videofile",
        "contentpart",
        "graphicframepr",
        "vba",
        "vbaproject",
    }
)
# These two auxiliary parts belong to the pinned python-pptx 1.0.2 template.
# They are not user images, rendered slides, or executable embedded packages.
_PRINTER_SETTINGS_SHA256 = "d7768f87e07d29634782e448ece1cddd05a52b9499254222b1190bc8f6dc579e"
_EXTENDED_PROPERTIES_MIME = "application/vnd.openxmlformats-officedocument.extended-properties+xml"
_AUXILIARY_HASHES = {
    "docProps/thumbnail.jpeg": "116f9aa8ea038e39ef47ee99a23592ab569cf77feafecd0a30e222369e931d25",
    "ppt/printerSettings/printerSettings1.bin": _PRINTER_SETTINGS_SHA256,
}
_FORBIDDEN_VISIBLE_TEXT = re.compile(
    r"\b(?:select|insert|update|delete|drop|alter|truncate|grant|revoke)\b|"
    r"\bcreate\s+(?:table|view|function|procedure)\b|"
    r"\b[a-z][a-z0-9+.-]*://|\bwww\.|\b(?:mailto|file|javascript):",
    re.IGNORECASE,
)


def _validate_visible_text(text: str) -> None:
    if _FORBIDDEN_VISIBLE_TEXT.search(text):
        raise PresentationArtifactInputError("Presentation contains prohibited visible text")


def _part_kind(member: str) -> str:
    if member.endswith(".rels"):
        return "relationships"
    fixed = {
        "docProps/core.xml": "core-properties",
        "docProps/app.xml": "extended-properties",
        "docProps/thumbnail.jpeg": "thumbnail",
        "ppt/presentation.xml": "officeDocument",
        "ppt/presProps.xml": "presProps",
        "ppt/viewProps.xml": "viewProps",
        "ppt/tableStyles.xml": "tableStyles",
        "ppt/printerSettings/printerSettings1.bin": "printerSettings",
    }
    if member in fixed:
        return fixed[member]
    for folder, kind in (
        ("slides", "slide"),
        ("slideLayouts", "slideLayout"),
        ("slideMasters", "slideMaster"),
        ("theme", "theme"),
    ):
        if member.startswith(f"ppt/{folder}/"):
            return kind
    raise PresentationArtifactInputError("Presentation OOXML unknown part kind")


def _expected_content_type(member: str) -> str:
    kind = _part_kind(member)
    fixed = {
        "relationships": "application/vnd.openxmlformats-package.relationships+xml",
        "core-properties": "application/vnd.openxmlformats-package.core-properties+xml",
        "extended-properties": _EXTENDED_PROPERTIES_MIME,
        "thumbnail": "image/jpeg",
        "officeDocument": _PPT_CT + "presentation.main+xml",
        "printerSettings": _PPT_CT + "printerSettings",
        "theme": "application/vnd.openxmlformats-officedocument.theme+xml",
    }
    return fixed.get(kind, _PPT_CT + kind + "+xml")


def _expected_xml_root(member: str) -> str:
    kind = _part_kind(member)
    if kind == "core-properties":
        return (
            "{http://schemas.openxmlformats.org/package/2006/metadata/core-properties}"
            "coreProperties"
        )
    if kind == "extended-properties":
        return (
            "{http://schemas.openxmlformats.org/officeDocument/2006/extended-properties}Properties"
        )
    roots = {
        "officeDocument": "presentation",
        "presProps": "presentationPr",
        "viewProps": "viewPr",
        "slide": "sld",
        "slideMaster": "sldMaster",
        "slideLayout": "sldLayout",
        "theme": "theme",
        "tableStyles": "tblStyleLst",
    }
    namespace = _DRAWING_NS if kind in {"theme", "tableStyles"} else _PRESENTATION_NS
    return f"{{{namespace}}}{roots[kind]}"


def _validate_content_types(root: ElementTree.Element, members: set[str]) -> None:
    if root.tag != f"{{{_CT_NS}}}Types":
        raise PresentationArtifactInputError("Presentation OOXML content types root invalid")
    defaults: dict[str, str] = {}
    overrides: dict[str, str] = {}
    allowed_defaults = {
        "xml": "application/xml",
        "rels": _expected_content_type("_rels/.rels"),
        "jpeg": "image/jpeg",
        "bin": _PPT_CT + "printerSettings",
    }
    for element in root:
        if element.tag == f"{{{_CT_NS}}}Default":
            extension = element.get("Extension", "")
            value = element.get("ContentType", "")
            if extension in defaults or value != allowed_defaults.get(extension):
                raise PresentationArtifactInputError(
                    "Presentation OOXML default content type invalid"
                )
            defaults[extension] = value
        elif element.tag == f"{{{_CT_NS}}}Override":
            name = element.get("PartName", "")
            member = name[1:] if name.startswith("/") else ""
            value = element.get("ContentType", "")
            if member not in members or member in overrides or member == "[Content_Types].xml":
                raise PresentationArtifactInputError(
                    "Presentation OOXML content type target invalid"
                )
            if value != _expected_content_type(member):
                raise PresentationArtifactInputError(
                    "Presentation OOXML override content type invalid"
                )
            overrides[member] = value
        else:
            raise PresentationArtifactInputError("Presentation OOXML content type element invalid")
    for member in members - {"[Content_Types].xml"}:
        declared = overrides.get(member, defaults.get(member.rsplit(".", 1)[-1], ""))
        if declared != _expected_content_type(member):
            raise PresentationArtifactInputError(
                "Presentation OOXML content type missing or inconsistent"
            )


def _validate_relationships(name: str, root: ElementTree.Element, members: set[str]) -> None:
    if root.tag != f"{{{_REL_NS}}}Relationships":
        raise PresentationArtifactInputError("Presentation OOXML relationship root invalid")
    base = str(PurePosixPath(name).parent.parent)
    if name != "_rels/.rels":
        source = posixpath.join(base, PurePosixPath(name).name.removesuffix(".rels"))
        if source not in members:
            raise PresentationArtifactInputError("Presentation OOXML relationship source missing")
    identifiers: set[str] = set()
    for relationship in root:
        identifier = relationship.get("Id", "")
        target = relationship.get("Target", "")
        if (
            relationship.tag != f"{{{_REL_NS}}}Relationship"
            or not identifier
            or identifier in identifiers
        ):
            raise PresentationArtifactInputError("Presentation OOXML relationship identity invalid")
        identifiers.add(identifier)
        if relationship.get("TargetMode", "Internal") != "Internal":
            raise PresentationArtifactInputError(
                "Presentation OOXML external relationship prohibited"
            )
        if not target or re.search(r"[%\\\x00-\x20]", target):
            raise PresentationArtifactInputError("Presentation OOXML relationship path invalid")
        uri = urlsplit(target)
        if uri.scheme or uri.netloc or uri.query or uri.fragment or target.startswith("/"):
            raise PresentationArtifactInputError("Presentation OOXML relationship URI prohibited")
        # Legitimate ../slideLayouts targets are resolved within the package.
        resolved = posixpath.normpath(posixpath.join(base, target))
        if (
            resolved.startswith("../")
            or resolved not in members
            or resolved == "[Content_Types].xml"
        ):
            raise PresentationArtifactInputError(
                "Presentation OOXML relationship target missing or unsafe"
            )
        kind = _part_kind(resolved)
        expected_type = (
            f"{_REL_NS}/metadata/{kind}"
            if kind in {"core-properties", "thumbnail"}
            else f"{_OFFICE_REL_NS}/{kind}"
        )
        if kind == "relationships" or relationship.get("Type") != expected_type:
            raise PresentationArtifactInputError("Presentation OOXML relationship type invalid")


def _validate_ooxml(content: bytes) -> None:
    """Validate generated OOXML without extraction or following relationships."""
    try:
        if len(content) > MAX_PRESENTATION_ARTIFACT_BYTES:
            raise PresentationArtifactInputError("Presentation OOXML byte limit exceeded")
        with ZipFile(BytesIO(content)) as archive:
            entries = archive.infolist()
            members = {entry.filename for entry in entries}
            if len(entries) > 256 or len(entries) != len(members):
                raise PresentationArtifactInputError("Presentation OOXML duplicate or excess parts")
            required = {"[Content_Types].xml", "_rels/.rels", "ppt/presentation.xml"}
            if not required.issubset(members):
                raise PresentationArtifactInputError("Presentation OOXML package is incomplete")
            if sum(entry.file_size for entry in entries) > 20 * 1024 * 1024:
                raise PresentationArtifactInputError(
                    "Presentation OOXML expanded byte limit exceeded"
                )
            roots: dict[str, ElementTree.Element] = {}
            for entry in entries:
                member = entry.filename
                if not _OOXML_MEMBER.fullmatch(member) or entry.flag_bits & 1:
                    raise PresentationArtifactInputError("Presentation OOXML unsupported member")
                data = archive.read(entry)
                if member.endswith((".xml", ".rels")):
                    if b"<!DOCTYPE" in data.upper() or b"<!ENTITY" in data.upper():
                        raise PresentationArtifactInputError(
                            "Presentation OOXML XML declarations prohibited"
                        )
                    # Only the UTF-8 XML emitted by the pinned generator is supported.
                    root = ElementTree.fromstring(data.decode("utf-8"))
                    roots[member] = root
                    for element in root.iter():
                        local_name = element.tag.rsplit("}", 1)[-1].casefold()
                        if local_name in _PROHIBITED_XML_ELEMENTS:
                            raise PresentationArtifactInputError(
                                "Presentation OOXML prohibited element"
                            )
                        if element.tag == f"{{{_DRAWING_NS}}}t":
                            _validate_visible_text(element.text or "")
                        if any(
                            key.rsplit("}", 1)[-1].casefold() == "action" for key in element.attrib
                        ):
                            raise PresentationArtifactInputError(
                                "Presentation OOXML actions prohibited"
                            )
                elif sha256(data).hexdigest() != _AUXILIARY_HASHES[member]:
                    raise PresentationArtifactInputError(
                        "Presentation OOXML auxiliary part changed"
                    )
            _validate_content_types(roots["[Content_Types].xml"], members)
            for member, root in roots.items():
                if member.endswith(".rels"):
                    _validate_relationships(member, root, members)
                elif member != "[Content_Types].xml" and root.tag != _expected_xml_root(member):
                    raise PresentationArtifactInputError("Presentation OOXML unexpected XML root")
    except (BadZipFile, ElementTree.ParseError, UnicodeError, ValueError, KeyError) as error:
        raise PresentationArtifactInputError(
            "Presentation OOXML package could not be validated"
        ) from error


class PresentationArtifactService:
    """Stateless generator for validated in-memory PowerPoint artifacts."""

    __slots__ = ()

    def build_pptx(
        self,
        result: AnalyticalPresentationResult,
    ) -> PresentationArtifact:
        """Build one safe PPTX without filesystem, database, LLM, or network access."""

        expected_slide_count = _validate_source(result)
        try:
            # A model_copy/model_construct object may bypass nested validators.
            # Rebuild from ordinary values, then render only the revalidated copy.
            result = AnalyticalPresentationResult.model_validate(result.model_dump(mode="python"))
        except ValidationError as error:
            raise PresentationArtifactInputError(
                "Presentation artifact nested source is invalid"
            ) from error
        artifact_id = _semantic_artifact_id(result)
        presentation = _build_presentation(result, artifact_id)

        if len(presentation.slides) != expected_slide_count:
            raise PresentationArtifactInputError(
                "Presentation artifact slide count is inconsistent"
            )

        buffer = BytesIO()
        presentation.save(buffer)
        content = buffer.getvalue()

        if len(content) > MAX_PRESENTATION_ARTIFACT_BYTES:
            raise PresentationArtifactSizeError("Presentation artifact exceeds the byte limit")

        _validate_ooxml(content)

        try:
            return PresentationArtifact(
                artifact_id=artifact_id,
                filename=f"analytical-presentation-{artifact_id}.pptx",
                content=content,
                size_bytes=len(content),
                slide_count=len(presentation.slides),
                source_row_count=result.source_row_count,
                presentation_version=result.presentation_version,
                visualization_version=result.visualizations.visualization_version,
                insight_version=result.insights.insight_version,
                analytics_version=result.visualizations.analytics_version,
                execution_version=result.visualizations.execution_version,
                semantic_version=result.visualizations.semantic_version,
                catalog_version=result.visualizations.catalog_version,
                media_type=PPTX_MIME_TYPE,
            )
        except ValidationError as error:
            raise PresentationArtifactInputError(
                "Presentation artifact contract is inconsistent"
            ) from error


def create_presentation_artifact_service() -> PresentationArtifactService:
    """Create one stateless presentation artifact generator."""

    return PresentationArtifactService()
