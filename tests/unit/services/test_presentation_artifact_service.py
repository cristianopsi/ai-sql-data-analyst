"""Tests for deterministic in-memory PowerPoint presentation artifacts."""

import inspect
import warnings
from decimal import Decimal
from io import BytesIO
from pathlib import Path
from typing import Literal
from xml.etree import ElementTree
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from pptx import Presentation
from pydantic import ValidationError

from backend.app.schemas.insights import (
    GroundedInsightClaim,
    GroundedInsightResult,
    InsightEvidenceReference,
    insight_claim_id,
)
from backend.app.schemas.llm import LLMTokenUsage
from backend.app.schemas.presentation import (
    AnalyticalPresentationResult,
    PresentationQueryResult,
)
from backend.app.schemas.presentation_artifact import (
    MAX_PRESENTATION_ARTIFACT_BYTES,
    MAX_PRESENTATION_SLIDES,
    PPTX_MIME_TYPE,
    PresentationArtifact,
)
from backend.app.schemas.visualization import (
    BarVisualizationItem,
    BarVisualizationSpec,
    DeterministicVisualizationResult,
    KPIVisualizationSpec,
    LineVisualizationPoint,
    LineVisualizationSpec,
    TableVisualizationRow,
    TableVisualizationSpec,
    visualization_specification_id,
)
from backend.app.services import presentation_artifact_service as artifact_service_module
from backend.app.services.presentation_artifact_service import (
    PresentationArtifactInputError,
    PresentationArtifactService,
    PresentationArtifactSizeError,
    create_presentation_artifact_service,
)


def _specification_id(
    chart_type: Literal["kpi", "table", "bar", "line"],
    metric: str,
    dimension: str | None = None,
) -> str:
    identifiers = (metric,) if dimension is None else (metric, dimension)
    return visualization_specification_id(chart_type, *identifiers)


def _presentation_result() -> AnalyticalPresentationResult:
    kpi = KPIVisualizationSpec.model_construct(
        spec_id=_specification_id("kpi", "approved_revenue"),
        chart_type="kpi",
        title="Approved revenue KPI",
        metric_name="approved_revenue",
        aggregation="sum",
        unit="brl",
        value_count=3,
        total=Decimal("300.00"),
        value=Decimal("300.00"),
        average=Decimal("100.00"),
        minimum=Decimal("50.00"),
        maximum=Decimal("150.00"),
    )
    table = TableVisualizationSpec.model_construct(
        spec_id=_specification_id("table", "approved_revenue", "region"),
        chart_type="table",
        title="Approved revenue by region",
        metric_name="approved_revenue",
        dimension_name="region",
        unit="brl",
        ranking_total=Decimal("300.00"),
        rows=(
            TableVisualizationRow.model_construct(
                position=1,
                label="North",
                value=Decimal("150.00"),
                share_percent=Decimal("50.00"),
            ),
            TableVisualizationRow.model_construct(
                position=2,
                label="South",
                value=Decimal("100.00"),
                share_percent=Decimal("33.3333"),
            ),
            TableVisualizationRow.model_construct(
                position=3,
                label="Central",
                value=Decimal("50.00"),
                share_percent=Decimal("16.6667"),
            ),
        ),
    )
    bar = BarVisualizationSpec.model_construct(
        spec_id=_specification_id("bar", "approved_revenue", "region"),
        chart_type="bar",
        title="Approved revenue by region",
        metric_name="approved_revenue",
        dimension_name="region",
        unit="brl",
        ranking_total=Decimal("300.00"),
        items=(
            BarVisualizationItem.model_construct(
                position=1,
                label="North",
                value=Decimal("150.00"),
                share_percent=Decimal("50.00"),
            ),
            BarVisualizationItem.model_construct(
                position=2,
                label="South",
                value=Decimal("100.00"),
                share_percent=Decimal("33.3333"),
            ),
            BarVisualizationItem.model_construct(
                position=3,
                label="Central",
                value=Decimal("50.00"),
                share_percent=Decimal("16.6667"),
            ),
        ),
    )
    line = LineVisualizationSpec.model_construct(
        spec_id=_specification_id("line", "approved_revenue", "month"),
        chart_type="line",
        title="Approved revenue trend",
        metric_name="approved_revenue",
        dimension_name="month",
        unit="brl",
        points=(
            LineVisualizationPoint.model_construct(
                position=1,
                label="2026-01",
                value=Decimal("50.00"),
                previous_value=None,
                absolute_change=None,
                percentage_change=None,
            ),
            LineVisualizationPoint.model_construct(
                position=2,
                label="2026-02",
                value=Decimal("100.00"),
                previous_value=Decimal("50.00"),
                absolute_change=Decimal("50.00"),
                percentage_change=Decimal("100.00"),
            ),
            LineVisualizationPoint.model_construct(
                position=3,
                label="2026-03",
                value=Decimal("150.00"),
                previous_value=Decimal("100.00"),
                absolute_change=Decimal("50.00"),
                percentage_change=Decimal("50.00"),
            ),
        ),
    )
    visualizations = DeterministicVisualizationResult.model_construct(
        visualization_version="1",
        visualization_status="specified",
        deterministic=True,
        analytics_version="1",
        execution_version="1",
        semantic_version="1",
        catalog_version="1",
        source_row_count=3,
        specifications=(kpi, table, bar, line),
    )
    evidence = (
        InsightEvidenceReference.model_construct(
            evidence_type="metric_summary",
            metric_name="approved_revenue",
            specification_id=None,
        ),
        InsightEvidenceReference.model_construct(
            evidence_type="visualization",
            metric_name=None,
            specification_id=bar.spec_id,
        ),
    )
    claim_text = "Approved revenue is led by the North region."
    claim = GroundedInsightClaim(
        claim_id=insight_claim_id(claim_text, evidence),
        text=claim_text,
        evidence=evidence,
    )
    insights = GroundedInsightResult.model_construct(
        insight_version="1",
        insight_status="generated",
        grounded=True,
        calculated_by_llm=False,
        analytics_version="1",
        visualization_version="1",
        execution_version="1",
        semantic_version="1",
        catalog_version="1",
        source_row_count=3,
        provider="mock",
        model="artifact-test-model",
        usage=LLMTokenUsage.model_construct(
            input_tokens=10,
            output_tokens=5,
            total_tokens=15,
        ),
        summary="Approved revenue is supported by deterministic evidence.",
        claims=(claim,),
    )
    query = PresentationQueryResult(
        validated_sql="SELECT raw_sql_marker FROM private_source",
        columns=("secret_column",),
        rows=(
            ("RAW-QUERY-ROW-ONE",),
            ("RAW-QUERY-ROW-TWO",),
            ("RAW-QUERY-ROW-THREE",),
        ),
        row_count=3,
    )
    result = AnalyticalPresentationResult(
        source_row_count=3,
        query=query,
        visualizations=visualizations,
        insights=insights,
    )
    return AnalyticalPresentationResult.model_validate(result.model_dump(mode="python"))


def _slide_text(artifact: PresentationArtifact) -> tuple[str, ...]:
    presentation = Presentation(BytesIO(artifact.content))
    texts: list[str] = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)

            if isinstance(text, str) and text:
                texts.append(text)

    return tuple(texts)


def _package_text(artifact: PresentationArtifact) -> bytes:
    with ZipFile(BytesIO(artifact.content)) as archive:
        return b"\n".join(
            archive.read(member)
            for member in archive.namelist()
            if member.endswith((".xml", ".rels"))
        )


def test_artifact_schema_is_strict_exact_and_immutable() -> None:
    artifact = PresentationArtifactService().build_pptx(_presentation_result())

    assert artifact.artifact_format == "pptx"
    assert artifact.media_type == PPTX_MIME_TYPE
    assert artifact.size_bytes == len(artifact.content)
    assert artifact.size_bytes <= MAX_PRESENTATION_ARTIFACT_BYTES
    assert artifact.slide_count <= MAX_PRESENTATION_SLIDES
    assert artifact.filename == f"analytical-presentation-{artifact.artifact_id}.pptx"

    with pytest.raises(ValidationError):
        PresentationArtifact.model_validate(
            {
                **artifact.model_dump(mode="python"),
                "filename": "../unsafe.pptx",
            }
        )

    with pytest.raises(ValidationError):
        artifact.filename = "changed.pptx"


def test_service_generates_openable_vector_only_pptx() -> None:
    artifact = PresentationArtifactService().build_pptx(_presentation_result())
    presentation = Presentation(BytesIO(artifact.content))

    assert artifact.content.startswith(b"PK\x03\x04")
    assert len(presentation.slides) == artifact.slide_count
    assert artifact.slide_count == 8

    with ZipFile(BytesIO(artifact.content)) as archive:
        members = tuple(archive.namelist())

    assert "ppt/presentation.xml" in members
    assert "ppt/vbaProject.bin" not in members
    assert not any(member.startswith("ppt/embeddings/") for member in members)
    assert not any(member.startswith("ppt/media/") for member in members)
    assert not any("oleObject" in member for member in members)


def test_service_is_semantically_repeatable_with_stable_identity() -> None:
    service = PresentationArtifactService()
    result = _presentation_result()

    first = service.build_pptx(result)
    second = service.build_pptx(result)

    assert first.artifact_id == second.artifact_id
    assert first.filename == second.filename
    assert first.slide_count == second.slide_count
    assert _slide_text(first) == _slide_text(second)


def test_service_uses_only_safe_projected_content() -> None:
    artifact = PresentationArtifactService().build_pptx(_presentation_result())
    package_text = _package_text(artifact)

    assert b"SELECT raw_sql_marker" not in package_text
    assert b"private_source" not in package_text
    assert b"secret_column" not in package_text
    assert b"RAW-QUERY-ROW" not in package_text
    assert b"artifact-test-model" not in package_text
    assert b"Approved revenue KPI" in package_text
    assert b"Approved revenue by region" in package_text
    assert b"Approved revenue trend" in package_text
    assert _presentation_result().insights.claims[0].claim_id.encode() in package_text

    source = inspect.getsource(artifact_service_module)
    assert "validated_sql" not in source
    assert "result.query" not in source


def test_service_preserves_evidence_and_source_versions() -> None:
    result = _presentation_result()
    artifact = PresentationArtifactService().build_pptx(result)
    rendered_text = "\n".join(_slide_text(artifact))
    claim = result.insights.claims[0]

    assert claim.claim_id in rendered_text
    assert claim.text in rendered_text
    assert "Metric evidence: approved_revenue" in rendered_text
    assert result.visualizations.specifications[2].spec_id in rendered_text
    assert "Presentation version" in rendered_text
    assert "Source row count" in rendered_text


def test_service_creates_no_filesystem_artifacts(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.chdir(tmp_path)

    artifact = PresentationArtifactService().build_pptx(_presentation_result())

    assert artifact.content
    assert tuple(tmp_path.iterdir()) == ()


def test_service_rejects_source_mismatch() -> None:
    result = _presentation_result().model_copy(
        update={
            "source_row_count": 4,
        }
    )

    with pytest.raises(
        PresentationArtifactInputError,
        match="row count",
    ):
        PresentationArtifactService().build_pptx(result)


def test_service_rejects_too_many_slides() -> None:
    result = _presentation_result()
    template = result.visualizations.specifications[0]
    specifications = tuple(
        template.model_copy(
            update={
                "spec_id": f"kpi-{'a' * 60}{index:04x}",
                "title": f"KPI {index}",
                "metric_name": f"metric_{index}",
            }
        )
        for index in range(17)
    )
    oversized_visualizations = result.visualizations.model_copy(
        update={
            "specifications": specifications,
        }
    )
    oversized = result.model_copy(
        update={
            "visualizations": oversized_visualizations,
        }
    )

    with pytest.raises(
        PresentationArtifactInputError,
        match="slide limit",
    ):
        PresentationArtifactService().build_pptx(oversized)


def test_service_fails_closed_for_oversized_bytes(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(
        artifact_service_module,
        "MAX_PRESENTATION_ARTIFACT_BYTES",
        100,
    )

    with pytest.raises(
        PresentationArtifactSizeError,
        match="byte limit",
    ):
        PresentationArtifactService().build_pptx(_presentation_result())


def test_factory_is_stateless() -> None:
    first = create_presentation_artifact_service()
    second = create_presentation_artifact_service()

    assert isinstance(first, PresentationArtifactService)
    assert isinstance(second, PresentationArtifactService)
    assert first is not second
    assert first.build_pptx(_presentation_result()).artifact_id == (
        second.build_pptx(_presentation_result()).artifact_id
    )


def test_d4_revalidates_nested_kpi_before_rendering(monkeypatch: pytest.MonkeyPatch) -> None:
    result = _presentation_result()
    specifications = result.visualizations.specifications
    forged = specifications[0].model_copy(update={"value": Decimal("999.00")})
    result = result.model_copy(
        update={
            "visualizations": result.visualizations.model_copy(
                update={"specifications": (forged, *specifications[1:])}
            )
        }
    )

    def forbidden_render(*args: object) -> None:
        pytest.fail("Invalid nested source reached the renderer")

    monkeypatch.setattr(artifact_service_module, "_build_presentation", forbidden_render)
    with pytest.raises(PresentationArtifactInputError, match="nested source"):
        PresentationArtifactService().build_pptx(result)


@pytest.mark.parametrize(
    "title", ["SELECT\tprivate_value FROM private_source", "https://example.invalid/d4"]
)
def test_d4_rejects_prohibited_visible_title(title: str) -> None:
    result = _presentation_result()
    specifications = result.visualizations.specifications
    changed = specifications[0].model_copy(update={"title": title})
    result = result.model_copy(
        update={
            "visualizations": result.visualizations.model_copy(
                update={"specifications": (changed, *specifications[1:])}
            )
        }
    )
    with pytest.raises(PresentationArtifactInputError, match="prohibited visible text"):
        PresentationArtifactService().build_pptx(result)


def _d4_mutated_package(content: bytes, mutation: str) -> bytes:
    """Inert fixtures only: never extract, execute, or fetch a package target."""
    with ZipFile(BytesIO(content)) as archive:
        parts = [(entry.filename, archive.read(entry)) for entry in archive.infolist()]
    slide = "ppt/slides/slide1.xml"
    rel = "ppt/slides/_rels/slide1.xml.rels"
    data = dict(parts)
    if mutation in {"unknown_member", "traversal_member", "backslash_member", "duplicate_member"}:
        name = {
            "unknown_member": "ppt/unexpected.bin",
            "traversal_member": "ppt/../probe.xml",
            "backslash_member": "ppt/..\\probe.xml",
            "duplicate_member": slide,
        }[mutation]
        parts.append((name, data[slide] if mutation == "duplicate_member" else b"inert"))
    else:
        target = slide
        replacement = data[slide]
        if mutation == "malformed_slide":
            replacement = b"<broken"
        elif mutation == "ole_element":
            root = ElementTree.fromstring(data[slide])
            ElementTree.SubElement(
                root, "{http://schemas.openxmlformats.org/presentationml/2006/main}oleObj"
            )
            replacement = ElementTree.tostring(root)
        elif mutation in {"macro_content_type", "duplicate_content_type"}:
            target = "[Content_Types].xml"
            root = ElementTree.fromstring(data[target])
            entry = next(item for item in root if item.get("PartName") == "/ppt/presentation.xml")
            if mutation == "macro_content_type":
                entry.set(
                    "ContentType",
                    "application/vnd.ms-powerpoint.presentation.macroEnabled.main+xml",
                )
            else:
                root.append(ElementTree.fromstring(ElementTree.tostring(entry)))
            replacement = ElementTree.tostring(root)
        elif mutation == "changed_auxiliary":
            target = "docProps/thumbnail.jpeg"
            replacement = b"inert altered template part"
        elif mutation == "doctype":
            replacement = b'<!DOCTYPE probe [<!ENTITY inert "value">]><probe/>'
        else:
            target = rel
            root = ElementTree.fromstring(data[rel])
            entry = root[0]
            destinations = {
                "external_without_mode": "https://example.invalid/inert",
                "traversal_target": "../../../../probe.xml",
                "dangling_target": "../slideLayouts/slideLayout999.xml",
                "encoded_target": "%2e%2e/slideLayouts/slideLayout1.xml",
                "fragment_target": "../slideLayouts/slideLayout1.xml#inert",
            }
            if mutation in destinations:
                entry.set("Target", destinations[mutation])
                entry.attrib.pop("TargetMode", None)
            elif mutation == "unknown_relationship_type":
                entry.set("Type", "urn:unsupported")
            elif mutation == "duplicate_relationship_id":
                root.append(ElementTree.fromstring(ElementTree.tostring(entry)))
            else:
                raise AssertionError("Unknown inert mutation")
            replacement = ElementTree.tostring(root)
        parts = [(name, replacement if name == target else value) for name, value in parts]
    buffer = BytesIO()
    with warnings.catch_warnings():
        warnings.filterwarnings("ignore", message="Duplicate name:.*", category=UserWarning)
        with ZipFile(buffer, "w", ZIP_DEFLATED) as archive:
            for name, value in parts:
                archive.writestr(name, value)
    return buffer.getvalue()


@pytest.mark.parametrize(
    "mutation",
    [
        "malformed_slide",
        "unknown_member",
        "traversal_member",
        "backslash_member",
        "duplicate_member",
        "external_without_mode",
        "traversal_target",
        "dangling_target",
        "ole_element",
        "macro_content_type",
        "duplicate_content_type",
        "changed_auxiliary",
        "doctype",
        "encoded_target",
        "fragment_target",
        "unknown_relationship_type",
        "duplicate_relationship_id",
    ],
)
def test_d4_rejects_inert_ooxml_mutations(mutation: str) -> None:
    artifact = PresentationArtifactService().build_pptx(_presentation_result())
    mutated = _d4_mutated_package(artifact.content, mutation)
    with pytest.raises(PresentationArtifactInputError):
        artifact_service_module._validate_ooxml(mutated)
